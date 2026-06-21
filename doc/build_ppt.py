#!/usr/bin/env python3
"""Build a Light-Professional themed deck for the bWAPP IT-audit presentation.
Navy header band, teal accent, severity colour coding, light panels, navy footer."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
NAVY  = RGBColor(0x14,0x2A,0x4F)
NAVY2 = RGBColor(0x1E,0x3A,0x5F)
TEAL  = RGBColor(0x12,0x8C,0x7E)
TEALL = RGBColor(0x53,0xC2,0xB2)
WHITE = RGBColor(0xFF,0xFF,0xFF)
PANEL = RGBColor(0xF1,0xF5,0xF9)
PANEL2= RGBColor(0xE5,0xEC,0xF3)
INK   = RGBColor(0x1B,0x25,0x33)
MUTED = RGBColor(0x5B,0x67,0x79)
LINE  = RGBColor(0xCF,0xD8,0xE2)
CRIT  = RGBColor(0xC0,0x39,0x2B)
HIGH  = RGBColor(0xD9,0x6B,0x12)
MED   = RGBColor(0xB8,0x86,0x09)
CODEBG= RGBColor(0x16,0x20,0x2C)
CODEFG= RGBColor(0xCF,0xD9,0xE0)
CODEGN= RGBColor(0x5BD,0x0,0x0) if False else RGBColor(0x6F,0xC2,0x76)
CODEYL= RGBColor(0xD7,0xBA,0x7D)

FONT="Calibri"; MONO="Consolas"
SEV={"CRITICAL":CRIT,"HIGH":HIGH,"MEDIUM":MED}

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
W=13.333

def _noline(sh): sh.line.fill.background()
def _fill(sh,c):
    sh.fill.solid(); sh.fill.fore_color.rgb=c

def rect(sl,x,y,w,h,color,rounded=False,line=None,lw=1.0):
    shp=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                            Inches(x),Inches(y),Inches(w),Inches(h))
    shp.shadow.inherit=False
    if color is None: shp.fill.background()
    else: _fill(shp,color)
    if line is None: _noline(shp)
    else:
        shp.line.color.rgb=line; shp.line.width=Pt(lw)
    if rounded:
        try: shp.adjustments[0]=0.08
        except Exception: pass
    return shp

def text(sl,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,
         space_after=2,line_spacing=1.0,wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,font,italic)."""
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf,m,0)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_before=Pt(0); p.space_after=Pt(space_after)
        p.line_spacing=line_spacing
        for (t,sz,col,bold,fn,it) in para:
            r=p.add_run(); r.text=t
            f=r.font; f.size=Pt(sz); f.bold=bold; f.italic=it; f.name=fn
            f.color.rgb=col
    return tb

def R(t,sz,col=INK,bold=False,fn=FONT,it=False): return (t,sz,col,bold,fn,it)

def footer(sl,n):
    rect(sl,0,7.18,W,0.32,NAVY)
    rect(sl,0,7.12,W,0.06,TEAL)
    text(sl,0.55,7.20,9,0.3,[[R("Audit Keamanan bWAPP   |   IT Audit   |   Ravi Arnan Irianto",9,RGBColor(0xC7,0xD2,0xDF))]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(sl,W-1.3,7.20,0.75,0.3,[[R(str(n),9,WHITE,True)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)

def base(kicker,title,n,sev=None):
    sl=prs.slides.add_slide(BLANK)
    rect(sl,0,0,W,1.18,NAVY)
    rect(sl,0,1.18,W,0.07,TEAL)
    text(sl,0.55,0.20,9,0.3,[[R(kicker.upper(),12,TEALL,True)]])
    text(sl,0.55,0.50,10.5,0.62,[[R(title,25,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    if sev:
        bw=2.0
        b=rect(sl,W-0.55-bw,0.38,bw,0.46,SEV[sev],rounded=True)
        text(sl,W-0.55-bw,0.38,bw,0.46,[[R(sev,13,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(sl,n)
    return sl

def bullets(sl,x,y,w,h,items,size=13,color=INK,gap=4,ls=1.05,marker=True,mcol=TEAL):
    paras=[]
    for it in items:
        if isinstance(it,tuple):
            lead,rest=it
            runs=[R(("•  " if marker else ""),size,mcol,True),R(lead+" ",size,INK,True),R(rest,size,color,False)]
        else:
            runs=[R(("•  " if marker else ""),size,mcol,True),R(it,size,color,False)]
        paras.append(runs)
    return text(sl,x,y,w,h,paras,space_after=gap,line_spacing=ls)

def code_block(sl,x,y,w,h,lines,title=None,size=10.5):
    if title:
        text(sl,x,y,w,0.3,[[R(title,11.5,TEAL,True)]]); y+=0.34; h-=0.34
    rect(sl,x,y,w,h,CODEBG,rounded=True)
    paras=[]
    for ln in lines:
        col=CODEFG
        s=ln
        if ln.startswith("$") or ln.startswith("#"): col=CODEGN
        paras.append([R(s if s else " ",size,col,False,MONO)])
    text(sl,x+0.18,y+0.14,w-0.36,h-0.28,paras,space_after=1,line_spacing=1.02)

def panel_title(sl,x,y,w,t,color=NAVY,size=13):
    text(sl,x,y,w,0.3,[[R(t,size,color,True)]])

def grid_table(sl,x,y,w,h,data,col_w,header=True,size=11,sev_idx=None,hsize=None):
    rows=len(data); cols=len(data[0])
    gt=sl.shapes.add_table(rows,cols,Inches(x),Inches(y),Inches(w),Inches(h)); tbl=gt.table
    # Table Grid style (thin borders), no auto banding
    sid=tbl._tbl.tblPr.find(qn('a:tableStyleId'))
    if sid is None:
        from pptx.oxml import parse_xml
        sid=tbl._tbl.tblPr.makeelement(qn('a:tableStyleId'),{}); tbl._tbl.tblPr.append(sid)
    sid.text='{5940675A-B579-460E-94D1-54222C63F5DA}'
    tbl.first_row=False; tbl.horz_banding=False
    for i,cw in enumerate(col_w): tbl.columns[i].width=Inches(cw)
    for r in range(rows):
        for c in range(cols):
            cell=tbl.cell(r,c)
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.10); cell.margin_right=Inches(0.08)
            cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            ishdr=(header and r==0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ishdr else (WHITE if r%2 else PANEL)
            tf=cell.text_frame; tf.word_wrap=True
            p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
            run=p.add_run(); run.text=str(data[r][c])
            f=run.font; f.name=FONT; f.size=Pt(hsize or size if ishdr else size)
            f.bold=ishdr
            if ishdr: f.color.rgb=WHITE
            else:
                f.color.rgb=INK
                if sev_idx is not None and c==sev_idx:
                    lvl=str(data[r][c]).upper()
                    if lvl in SEV: f.color.rgb=SEV[lvl]; f.bold=True
    return tbl

# ============================================================ SLIDE 1 : TITLE
sl=prs.slides.add_slide(BLANK)
rect(sl,0,0,W,7.5,WHITE)
rect(sl,0,0,W,3.05,NAVY)
rect(sl,0,3.05,W,0.10,TEAL)
text(sl,0.85,0.62,11,0.4,[[R("VULNERABILITY ASSESSMENT",14,TEALL,True)]])
text(sl,0.85,1.05,11.6,1.7,[[R("LAPORAN AUDIT KEAMANAN",37,WHITE,True)],
                            [R("APLIKASI WEB",37,WHITE,True)]],line_spacing=1.02)
text(sl,0.88,3.45,11,0.5,[[R("Vulnerability Assessment terhadap bWAPP (Buggy Web Application)",16,MUTED)]])
# metadata panel
rect(sl,0.85,4.25,8.4,2.35,PANEL,rounded=True)
rect(sl,0.85,4.25,0.12,2.35,TEAL)
meta=[("Auditor","Ravi Arnan Irianto"),("Mata Kuliah","IT Audit"),
      ("Target","bWAPP (OWASP)  -  Docker, http://localhost:8080"),
      ("Standar","OWASP Top 10 2021")]
mp=[]
for k,v in meta: mp.append([R(k+":  ",14,NAVY,True),R(v,14,INK)])
text(sl,1.25,4.5,7.8,1.9,mp,space_after=9,line_spacing=1.0,anchor=MSO_ANCHOR.MIDDLE)
try:
    sl.shapes.add_picture("unud_logo.png",Inches(10.4),Inches(4.35),height=Inches(2.05))
except Exception as e: print("logo:",e)

# ============================================================ SLIDE 2 : APA ITU bWAPP
sl=base("01  Pengenalan","Apa itu bWAPP?",2)
text(sl,0.55,1.55,12.2,1.1,[[R("bWAPP (Buggy Web Application)",14,INK,True),
   R(" adalah aplikasi web PHP/MySQL yang ",14,INK),R("sengaja dibuat rentan",14,CRIT,True),
   R(" untuk tujuan edukasi keamanan, dikembangkan oleh Malik Mesellem. Aplikasi ini menjadi sarana legal untuk mempraktikkan teknik audit dan pengujian penetrasi dalam lingkungan terisolasi.",14,INK)]],
   line_spacing=1.12)
cards=[("100+","Jenis kerentanan tersistematis"),("3","Tingkat kesulitan (Low / Medium / High)"),
       ("2","Metode deploy (Docker / Bee-Box VM)")]
cx=0.55; cw=4.0; gap=0.18
for i,(big,lab) in enumerate(cards):
    x=cx+i*(cw+gap)
    rect(sl,x,3.05,cw,1.85,PANEL,rounded=True)
    rect(sl,x,3.05,cw,0.12,TEAL,rounded=False)
    text(sl,x,3.35,cw,0.95,[[R(big,46,NAVY,True)]],align=PP_ALIGN.CENTER)
    text(sl,x+0.25,4.32,cw-0.5,0.5,[[R(lab,12.5,MUTED)]],align=PP_ALIGN.CENTER)
rect(sl,0.55,5.25,12.23,0.9,PANEL2,rounded=True)
rect(sl,0.55,5.25,0.12,0.9,TEAL)
text(sl,0.9,5.25,12.0,0.9,[[R("Catatan:  ",13,NAVY,True),
   R("bWAPP bukan aplikasi yang error, melainkan dirancang untuk dieksploitasi sebagai media latihan.",13,INK)]],
   anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ SLIDE 3 : RELEVANSI
sl=base("02  Relevansi","Kenapa bWAPP untuk IT Audit?",3)
bullets(sl,0.55,1.7,7.0,3.5,[
  ("Coverage terluas:","100+ kerentanan dalam satu platform."),
  ("Selaras OWASP:","tiap bug mudah dipetakan ke framework standar."),
  ("Setup cepat:","siap pakai via Docker dalam hitungan menit."),
  ("Bertingkat:","3 level kesulitan untuk simulasi bertahap."),
],size=14,gap=12,ls=1.1)
grid_table(sl,7.95,1.75,4.85,2.4,
  [["Platform","Jumlah Bug"],["bWAPP","100+"],["Juice Shop","~75"],["WebGoat","~35"],["DVWA","~15"]],
  [3.0,1.85],size=12.5)
text(sl,7.95,4.45,4.85,0.8,[[R("bWAPP mengungguli platform sejenis dalam cakupan kerentanan.",12,MUTED,it=True)]],line_spacing=1.05)

# ============================================================ SLIDE 4 : LINGKUNGAN
sl=base("03  Lingkungan","Setup dan Lingkungan Lab",4)
text(sl,0.55,1.55,6.0,0.35,[[R("Deployment via Docker (terisolasi ke localhost):",13,NAVY,True)]])
code_block(sl,0.55,1.95,6.0,1.5,[
  "$ docker run -d -p 8080:80 \\",
  "    hackersploit/bwapp-docker",
  "# inisialisasi DB sekali:",
  "  http://localhost:8080/install.php"],size=11)
bullets(sl,0.55,3.7,6.0,2.2,[
  ("Akses:","http://localhost:8080"),
  ("Login:","bee / bug (kredensial lab default)"),
  ("Security level:","Low (security_level = 0)"),
  ("Stack:","PHP 5.5.9, Apache 2.4.7, MySQL, Ubuntu"),
],size=12.5,gap=7)
rect(sl,6.9,1.6,5.9,4.45,PANEL,rounded=True)
panel_title(sl,7.2,1.85,5.4,"Ruang Lingkup Audit")
bullets(sl,7.2,2.35,5.35,3.4,[
  ("Jenis uji:","black-box dan grey-box"),
  ("In-scope:","web app layer (injection, XSS, info disclosure, misconfig, header)"),
  ("Out-of-scope:","audit OS/host, DoS, social engineering, sistem nyata"),
  ("Sifat:","lab terisolasi, aplikasi sengaja rentan"),
],size=12.5,gap=9,ls=1.08)

# ============================================================ SLIDE 5 : 5 TAHAP
sl=base("04  Metodologi","Alur Audit: 5 Tahap",5)
steps=[("1","Setup","Deploy bWAPP via Docker dan verifikasi akses."),
       ("2","Scope","Pilih kategori kerentanan dan tetapkan level kesulitan."),
       ("3","Pengujian","Scan otomatis, eksploitasi terarah, verifikasi manual."),
       ("4","Analisis & Mapping","Petakan ke OWASP Top 10 2021, tetapkan severity."),
       ("5","Pelaporan","Susun temuan, bukti, dan rekomendasi remediasi.")]
cw=2.38; gap=0.10; x0=0.55; y0=2.1
for i,(num,tt,dd) in enumerate(steps):
    x=x0+i*(cw+gap)
    rect(sl,x,y0,cw,3.3,PANEL,rounded=True)
    circ=sl.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+cw/2-0.45),Inches(y0+0.28),Inches(0.9),Inches(0.9))
    circ.shadow.inherit=False; _fill(circ,NAVY); _noline(circ)
    text(sl,x+cw/2-0.45,y0+0.28,0.9,0.9,[[R(num,26,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(sl,x+0.12,y0+1.35,cw-0.24,0.7,[[R(tt,13.5,NAVY,True)]],align=PP_ALIGN.CENTER,line_spacing=1.0)
    text(sl,x+0.18,y0+2.0,cw-0.36,1.2,[[R(dd,11.5,MUTED)]],align=PP_ALIGN.CENTER,line_spacing=1.05)
    if i<4:
        text(sl,x+cw-0.02,y0+0.55,0.3,0.6,[[R("›",24,TEAL,True)]],align=PP_ALIGN.CENTER)

# ============================================================ SLIDE 6 : TOOLS
sl=base("04  Metodologi","Tools yang Digunakan",6)
grid_table(sl,0.55,1.6,12.23,3.7,[
  ["Tool","Versi","Fungsi"],
  ["Nikto","2.1.5","Web server scanner (info disclosure, misconfiguration)"],
  ["SQLMap","1.10.2.6","Deteksi dan eksploitasi SQL Injection otomatis"],
  ["OWASP ZAP","2.17.0","Dynamic Application Security Testing (DAST)"],
  ["Retire.js","5.4.3","Deteksi pustaka JavaScript usang (SCA)"],
  ["Nmap + vulners","7.94SVN","Deteksi versi layanan dan pemetaan CVE server"],
  ["Manual (curl, browser)","-","Verifikasi manual XSS dan kebocoran data"],
],[2.7,1.7,7.83],size=12.5)
rect(sl,0.55,5.55,12.23,0.85,PANEL2,rounded=True)
rect(sl,0.55,5.55,0.12,0.85,TEAL)
text(sl,0.9,5.55,12.0,0.85,[[R("Kombinasi scan otomatis + eksploitasi terarah + analisis komponen + uji manual untuk menutup blind spot tiap tool.",12.5,INK)]],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.05)

# ============================================================ SLIDE 7 : RINGKASAN EKSEKUTIF
sl=base("05  Hasil","Ringkasan Eksekutif",7)
text(sl,0.55,1.6,7.4,3.8,[
  [R("Audit menemukan ",14,INK),R("5 area temuan utama",14,NAVY,True),R(" dengan severity Medium hingga Critical.",14,INK)],
  [R("Dua temuan ",14,INK),R("Critical",14,CRIT,True),R(" memungkinkan penyerang mengambil alih data sensitif tanpa autentikasi: kebocoran kredensial melalui file konfigurasi yang terekspos, dan SQL Injection yang membongkar seluruh tabel pengguna termasuk akun admin.",14,INK)],
  [R("Aplikasi juga rentan XSS, memakai komponen End-of-Life, dan memiliki sejumlah security misconfiguration.",14,INK)],
],space_after=10,line_spacing=1.15)
# severity distribution panel
rect(sl,8.35,1.6,4.43,4.55,PANEL,rounded=True)
text(sl,8.35,1.85,4.43,0.4,[[R("Distribusi Severity",13.5,NAVY,True)]],align=PP_ALIGN.CENTER)
dist=[("Critical","2",CRIT),("High","2",HIGH),("Medium","1",MED)]
yy=2.45
for nm,ct,col in dist:
    rect(sl,8.7,yy,3.73,0.78,WHITE,rounded=True,line=LINE,lw=0.75)
    rect(sl,8.7,yy,0.14,0.78,col)
    text(sl,9.0,yy,2.4,0.78,[[R(nm,14,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    text(sl,11.0,yy,1.25,0.78,[[R(ct,24,col,True)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    yy+=0.92
text(sl,8.7,yy+0.05,3.73,0.7,[[R("Total 5 temuan utama",12,MUTED)],
     [R("(plus 6 temuan Low pendukung)",12,MUTED)]],align=PP_ALIGN.CENTER,line_spacing=1.05)

# ============================================================ SLIDE 8 : RINGKASAN TEMUAN
sl=base("05  Hasil","Ringkasan Temuan",8)
grid_table(sl,0.55,1.7,12.23,4.2,[
  ["ID","Temuan","Severity","OWASP 2021","Tool"],
  ["F-01","Sensitive Data Exposure (kredensial bocor)","Critical","A01, A02, A05","Nikto + manual"],
  ["F-02","SQL Injection","Critical","A03, A02","SQLMap"],
  ["F-03","Cross-Site Scripting (Reflected + Stored)","High","A03","Manual"],
  ["F-04","Security Misconfiguration & Missing Headers","Medium","A05, A01","OWASP ZAP"],
  ["F-05","Vulnerable & Outdated Components","High","A06","Retire.js + Nmap"],
],[1.0,5.0,1.5,2.5,2.23],size=12.5,sev_idx=2)

# ============================================================ FINDING SLIDES 9-13
def finding(n,fid,ftitle,sev,owasp,desc,bukti_title,code_lines,dampak,rekom):
    sl=base("06  Temuan Terperinci",f"{fid}   {ftitle}",n,sev=sev)
    text(sl,0.55,1.42,8.0,0.32,[[R("OWASP 2021:  ",11.5,TEAL,True),R(owasp,11.5,MUTED)]])
    # left column
    text(sl,0.55,1.95,6.05,1.5,[[R(desc,12.5,INK)]],line_spacing=1.12)
    rect(sl,0.55,3.45,6.05,0.86,RGBColor(0xFB,0xEC,0xEC) if sev in("CRITICAL","HIGH") else RGBColor(0xFB,0xF4,0xDF),rounded=True)
    rect(sl,0.55,3.45,0.12,0.86,SEV[sev])
    text(sl,0.9,3.45,5.6,0.86,[[R("Dampak:  ",12.5,SEV[sev],True),R(dampak,12,INK)]],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.05)
    panel_title(sl,0.55,4.5,6.0,"Rekomendasi Remediasi",color=NAVY,size=12.5)
    bullets(sl,0.55,4.9,6.05,2.1,rekom,size=11.5,gap=4,ls=1.05)
    # right column code
    code_block(sl,6.85,1.95,5.95,4.8,code_lines,title=bukti_title,size=10)
    return sl

finding(9,"F-01","Sensitive Data Exposure","CRITICAL",
  "A01 Broken Access Control, A02 Cryptographic Failures, A05 Security Misconfiguration",
  "File konfigurasi dan backup berisi kredensial dapat diakses publik tanpa autentikasi. Directory indexing yang aktif menampilkan isi folder sensitif seperti /passwords/.",
  "Bukti: config.inc terekspos",
  ['$ curl http://localhost:8080/config.inc','$server   = "localhost";','$username = "bwapp";',
   '$password = "bwApped";','$database = "bWAPP";','','/passwords/ -> heroes.xml,','   web.config.bak, wp-config.bak','   (password plaintext)'],
  "kredensial DB diperoleh tanpa eksploitasi, berpotensi full database compromise.",
  ["Matikan directory indexing (Options -Indexes).","Hapus file backup (.bak, .inc, .old) dari web root.",
   "Pindahkan file konfigurasi ke luar document root.","Jangan simpan password plaintext, gunakan hashing kuat.",
   "Hapus installer (install.php) setelah instalasi."])

finding(10,"F-02","SQL Injection","CRITICAL",
  "A03 Injection, A02 Cryptographic Failures",
  "Parameter title pada modul SQL Injection (GET/Search) tidak disanitasi, sehingga input pengguna dieksekusi langsung sebagai bagian dari query SQL.",
  "Bukti: SQLMap dump bWAPP.users",
  ["Parameter 'title' rentan 4 teknik:","boolean-blind, error, time-blind, UNION","",
   "Database: bWAPP  Table: users","id | login  | admin | password(SHA1)","1  | A.I.M. |  1    | 6885858486..affd0",
   "2  | bee    |  1    | 6885858486..affd0",'echo -n "bug" | sha1sum -> match','Password admin dipulihkan: "bug"'],
  "akses baca penuh ke database; hash SHA1 unsalted mudah dipecahkan.",
  ["Gunakan prepared statements / parameterized query (PDO, MySQLi).","Terapkan input validation berbasis whitelist.",
   "Ganti hashing password ke bcrypt atau argon2 dengan salt unik.","Terapkan least privilege pada akun database aplikasi."])

finding(11,"F-03","Cross-Site Scripting (XSS)","HIGH",
  "A03 Injection (Cross-Site Scripting)",
  "Aplikasi tidak melakukan output encoding maupun sanitasi input. Ditemukan dua varian: Reflected (GET) dan Stored (Blog) yang tersimpan permanen.",
  "Bukti: payload tereksekusi",
  ["Reflected (xss_get.php, firstname):","Welcome <script>alert('XSS-bWAPP')","        </script> test","",
   "Stored (xss_stored_1.php):","<td><script>alert('Stored-XSS-","     by-Audit')</script></td>","",
   "(payload stored telah dibersihkan)"],
  "pencurian session cookie, session hijacking, defacement, phishing. Diperparah cookie tanpa HttpOnly.",
  ["Terapkan output encoding sesuai konteks (htmlspecialchars).","Lakukan input validation di sisi server.",
   "Pasang Content Security Policy (CSP).","Set flag HttpOnly dan Secure pada cookie session."])

finding(12,"F-04","Security Misconfiguration","MEDIUM",
  "A05 Security Misconfiguration, A01 Broken Access Control",
  "OWASP ZAP menemukan sejumlah kesalahan konfigurasi dan header keamanan yang hilang: total 5 temuan Medium dan 6 temuan Low pendukung.",
  "Bukti: temuan Medium dari ZAP",
  ["- Absence of Anti-CSRF Tokens   (CWE-352)","- Application Error Disclosure  (CWE-550)",
   "- CSP Header Not Set            (CWE-693)","- Directory Browsing            (CWE-548)",
   "- Missing Anti-clickjacking     (CWE-1021)","","Low: cookie tanpa HttpOnly/SameSite,",
   "banner Apache/2.4.7 & PHP/5.5.9 bocor,","X-Content-Type-Options tak diset"],
  "memperbesar peluang CSRF, clickjacking, dan info disclosure bagi penyerang.",
  ["Tambahkan token Anti-CSRF unik pada setiap form pengubah state.","Set header CSP, X-Frame-Options, X-Content-Type-Options: nosniff.",
   "Set flag HttpOnly, Secure, SameSite pada cookie.","Gunakan custom error page; sembunyikan banner versi.",
   "Update komponen EOL (PHP 5.5.9, Apache 2.4.7)."])

finding(13,"F-05","Vulnerable & Outdated Components","HIGH",
  "A06 Vulnerable and Outdated Components",
  "Aplikasi berjalan di atas komponen End-of-Life pada sisi klien (jQuery) dan server (Apache, PHP) yang tidak lagi menerima patch keamanan.",
  "Bukti: Retire.js + Nmap vulners",
  ["jQuery 1.4.4 (2010, EOL) -> 7 CVE:","CVE-2020-11023, CVE-2019-11358,","CVE-2020-7656, CVE-2015-9251, dll","",
   "Apache httpd 2.4.7 (Ubuntu):","vulners cocokkan 106 CVE unik","(Critical 24 | High 42 | Medium 40)","",
   "PHP 5.5.9 -> EOL sejak Juli 2016"],
  "permukaan serangan luas (XSS, prototype pollution, request smuggling) dengan banyak exploit publik.",
  ["Perbarui jQuery ke versi didukung (3.5.0+).","Update Apache; migrasikan PHP 5.5.9 ke PHP 8.x.",
   "Terapkan patch management dan SCA berkala.","Skor per-CVE potensial (berbasis banner); fakta EOL terkonfirmasi High."])

# ============================================================ SLIDE 14 : CATATAN METODOLOGI
sl=base("07  Catatan Metodologi","Keterbatasan Automated Scanner",14)
rect(sl,0.55,1.6,12.23,1.0,RGBColor(0xFB,0xF4,0xDF),rounded=True)
rect(sl,0.55,1.6,0.12,1.0,HIGH)
text(sl,0.9,1.6,11.7,1.0,[[R("OWASP ZAP melaporkan 0 temuan High",13.5,HIGH,True),
   R(", padahal pengujian manual dan SQLMap membuktikan adanya SQL Injection (Critical) dan Stored XSS (High) pada aplikasi yang sama.",13.5,INK)]],
   anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.08)
text(sl,0.55,2.85,12.23,1.3,[[R("Penyebab:  ",13,NAVY,True),
   R("ZAP hanya meng-crawl halaman permukaan dan tidak masuk ke modul kerentanan di balik autentikasi atau yang tidak ter-link otomatis. Automated DAST scanner memiliki blind spot, sehingga audit yang baik wajib mengombinasikan scan otomatis, eksploitasi terarah, dan uji manual.",13,INK)]],line_spacing=1.15)
panel_title(sl,0.55,4.35,12,"Cross-Validation  (dikonfirmasi 2 tool berbeda: Nikto + ZAP)",color=NAVY,size=13)
bullets(sl,0.55,4.8,12.0,1.8,[
  "Directory Browsing aktif","Kebocoran banner versi server (Apache/2.4.7)",
  "Missing X-Frame-Options (anti-clickjacking)","Cookie session tanpa flag HttpOnly"],size=13,gap=6)

# ============================================================ SLIDE 15 : ETIKA LAB
sl=base("08  Etika Lab","Best Practice dan Keamanan Lab",15)
rect(sl,0.55,1.65,6.0,4.7,RGBColor(0xEC,0xF6,0xF1),rounded=True)
rect(sl,0.55,1.65,6.0,0.6,TEAL,rounded=True)
text(sl,0.55,1.65,6.0,0.6,[[R("✓  YANG HARUS DILAKUKAN",14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
bullets(sl,0.85,2.5,5.45,3.7,[
  "Bind Docker hanya ke localhost.","Gunakan Host-only Adapter pada VM.","Ambil snapshot sebelum eksploitasi.",
  "Dokumentasikan setiap langkah (payload, screenshot, respons).","Bersihkan payload stored setelah verifikasi."],
  size=12.5,gap=9,ls=1.08,mcol=TEAL)
rect(sl,6.8,1.65,6.0,4.7,RGBColor(0xFB,0xEC,0xEC),rounded=True)
rect(sl,6.8,1.65,6.0,0.6,CRIT,rounded=True)
text(sl,6.8,1.65,6.0,0.6,[[R("✗  YANG TIDAK BOLEH",14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
bullets(sl,7.1,2.5,5.45,3.7,[
  "Expose port ke jaringan publik.","Menjalankan di komputer produksi.",
  "Menggunakan payload terhadap target nyata tanpa izin tertulis (ilegal).",
  "Memperlakukan data dump sebagai kredensial nyata (ini user seeded bWAPP)."],
  size=12.5,gap=9,ls=1.08,mcol=CRIT)

# ============================================================ SLIDE 16 : KESIMPULAN
sl=base("09  Kesimpulan","Kesimpulan",16)
concl=[("Kerentanan serius dan beragam.","bWAPP mencakup mayoritas kategori OWASP Top 10 2021."),
  ("Dua Critical jadi prioritas.","Sensitive Data Exposure dan SQL Injection memungkinkan kompromi data langsung."),
  ("Komponen usang memperluas serangan.","A06: jQuery 1.4.4, Apache 2.4.7, dan PHP 5.5.9 sudah End-of-Life."),
  ("Kombinasi metode terbukti penting.","Tidak ada satu tool pun yang menemukan seluruh kerentanan."),
  ("Media latihan audit yang efektif.","Dari identifikasi, eksploitasi, hingga rekomendasi selaras standar industri.")]
yy=1.75
for lead,rest in concl:
    rect(sl,0.55,yy,12.23,0.84,PANEL,rounded=True)
    rect(sl,0.55,yy,0.12,0.84,TEAL)
    text(sl,0.95,yy,11.7,0.84,[[R(lead+"  ",13.5,NAVY,True),R(rest,13,INK)]],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.05)
    yy+=0.96

# ============================================================ SLIDE 17 : THANKS
sl=prs.slides.add_slide(BLANK)
rect(sl,0,0,W,7.5,NAVY)
rect(sl,0.85,2.45,3.2,0.12,TEAL)
text(sl,0.85,1.5,11.6,1.0,[[R("Terima Kasih",46,WHITE,True)]])
text(sl,0.88,2.7,11.6,0.55,[[R("Ada pertanyaan?",20,TEALL)]])
rect(sl,0.85,3.85,8.9,2.4,NAVY2,rounded=True)
rect(sl,0.85,3.85,0.12,2.4,TEAL)
text(sl,1.2,4.1,8.3,1.95,[
  [R("Auditor:  ",14,TEALL,True),R("Ravi Arnan Irianto  -  Mata Kuliah IT Audit",14,WHITE)],
  [R("Demo cepat:  ",14,TEALL,True),R("http://localhost:8080  (login bee / bug)",14,WHITE)],
  [R("Referensi:  ",14,TEALL,True),R("owasp.org/Top10  -  itsecgames.com",14,WHITE)],
],space_after=12,line_spacing=1.0,anchor=MSO_ANCHOR.MIDDLE)

prs.save("presentasi_audit_bwapp.pptx")
print("saved presentasi_audit_bwapp.pptx with",len(prs.slides.__iter__.__self__._sldIdLst),"slides")
