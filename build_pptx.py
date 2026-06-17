#!/usr/bin/env python3
"""
Generator slide presentasi hasil audit keamanan bWAPP.
Output: doc/presentasi_audit_bwapp.pptx (16:9).

Deck gabungan: pengenalan bWAPP -> metodologi -> hasil audit (F-01..F-05)
-> catatan metodologi -> etika lab -> kesimpulan. Tema navy, konsisten
dengan laporan_audit_akhir.html. Tanpa emoji, tanpa em dash.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
NAVY   = RGBColor(0x00, 0x33, 0x66)
NAVY2  = RGBColor(0x14, 0x47, 0x7D)
CRIT   = RGBColor(0xB3, 0x00, 0x00)
HIGH   = RGBColor(0xE8, 0x62, 0x0C)
MED    = RGBColor(0xE0, 0xA8, 0x00)
LOW    = RGBColor(0xB8, 0xB0, 0x00)
LIGHT  = RGBColor(0xF3, 0xF5, 0xF8)
PANEL  = RGBColor(0xEE, 0xF2, 0xF7)
TEXT   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x88, 0x88, 0x88)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
CODEBG = RGBColor(0x1E, 0x1E, 0x1E)
CODEFG = RGBColor(0xD4, 0xD4, 0xD4)

FONT = "Calibri"
MONO = "Consolas"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs, each a list of (text, size, color, bold, italic, font)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, italic, font) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
    return tb


def R(t, size=14, color=TEXT, bold=False, italic=False, font=FONT):
    return (t, size, color, bold, italic, font)


def bullets(s, x, y, w, h, items, size=15, gap=7, color=TEXT, marker="-  "):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        if isinstance(it, tuple):
            label, rest = it
            r = p.add_run(); r.text = marker + label
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = NAVY2; r.font.name = FONT
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = FONT
        else:
            r = p.add_run(); r.text = marker + it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(s, kicker, title, num):
    """Standard content-slide header: top navy bar + kicker + title + page no."""
    rect(s, 0, 0, EMU_W, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), EMU_W, Inches(0.06), ACCENT)
    if kicker:
        txt(s, Inches(0.6), Inches(0.16), Inches(11), Inches(0.3),
            [[R(kicker, 12, RGBColor(0x9F, 0xC5, 0xE8), True, False, FONT)]])
    txt(s, Inches(0.6), Inches(0.40), Inches(11.5), Inches(0.7),
        [[R(title, 26, WHITE, True, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(12.4), Inches(0.16), Inches(0.7), Inches(0.3),
        [[R(num, 11, RGBColor(0x9F, 0xC5, 0xE8), False, False, FONT)]],
        align=PP_ALIGN.RIGHT)


def footer(s):
    txt(s, Inches(0.6), Inches(7.12), Inches(9), Inches(0.3),
        [[R("Audit Keamanan bWAPP  |  IT Audit  |  Ravi Arnan Irianto",
            9, GRAY, False, False, FONT)]])


def chip(s, x, y, label, color, w=Inches(1.35), h=Inches(0.38)):
    sp = rect(s, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(13); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = WHITE if color in (CRIT, HIGH, NAVY) else TEXT
    return sp


def codebox(s, x, y, w, h, lines, size=11.5):
    rect(s, x, y, w, h, CODEBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(10); tf.margin_right = Pt(8)
    tf.margin_top = Pt(6); tf.margin_bottom = Pt(6)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05; p.space_after = Pt(1)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.name = MONO; r.font.color.rgb = CODEFG
    return tb


def set_cell(cell, text, size=12, color=TEXT, bold=False, fill=None,
             align=PP_ALIGN.LEFT, font=FONT):
    cell.margin_left = Pt(7); cell.margin_right = Pt(7)
    cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.name = font


def table(s, x, y, w, rows, cols, col_widths, row_height=Inches(0.4)):
    gt = s.shapes.add_table(rows, cols, x, y, w, row_height * rows).table
    gt.first_row = False; gt.horz_banding = False
    for j, cw in enumerate(col_widths):
        gt.columns[j].width = cw
    # strip default table style border noise -> keep simple
    return gt


def sev_color(sev):
    return {"Critical": CRIT, "High": HIGH, "Medium": MED, "Low": LOW}[sev]


# ================================================================ SLIDE 1: cover
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, NAVY)
rect(s, 0, Inches(4.95), EMU_W, Inches(0.07), ACCENT)
# subtle accent block
rect(s, Inches(0.0), Inches(0.0), Inches(0.35), EMU_H, ACCENT)
txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(2.2),
    [[R("LAPORAN AUDIT KEAMANAN", 40, WHITE, True, False, FONT)],
     [R("APLIKASI WEB", 40, WHITE, True, False, FONT)]],
    line_spacing=1.05)
txt(s, Inches(0.9), Inches(3.7), Inches(11), Inches(0.9),
    [[R("Vulnerability Assessment terhadap bWAPP (Buggy Web Application)",
        20, RGBColor(0xCF, 0xE0, 0xF0), False, False, FONT)]])
txt(s, Inches(0.9), Inches(5.25), Inches(11), Inches(1.6),
    [[R("Auditor:  ", 16, GRAY, True, False, FONT), R("Ravi Arnan Irianto", 16, WHITE, False, False, FONT)],
     [R("Mata Kuliah:  ", 16, GRAY, True, False, FONT), R("IT Audit", 16, WHITE, False, False, FONT)],
     [R("Target:  ", 16, GRAY, True, False, FONT), R("bWAPP (OWASP)  -  Docker, http://localhost:8080", 16, WHITE, False, False, FONT)],
     [R("Standar:  ", 16, GRAY, True, False, FONT), R("OWASP Top 10 2021", 16, WHITE, False, False, FONT)]],
    space_after=6)


# ================================================================ SLIDE 2: apa itu bWAPP
s = slide()
header(s, "01  PENGENALAN", "Apa itu bWAPP?", "2")
txt(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(1.1),
    [[R("bWAPP (Buggy Web Application) adalah aplikasi web PHP/MySQL yang ",
        16, TEXT, False, False, FONT),
      R("sengaja dibuat rentan", 16, CRIT, True, False, FONT),
      R(" untuk tujuan edukasi keamanan, dikembangkan oleh Malik Mesellem. "
        "Aplikasi ini menjadi sarana legal untuk mempraktikkan teknik audit "
        "dan pengujian penetrasi dalam lingkungan terisolasi.",
        16, TEXT, False, False, FONT)]], line_spacing=1.1)
# three stat cards
cards = [("100+", "Jenis kerentanan tersistematis"),
         ("3", "Tingkat kesulitan (Low / Medium / High)"),
         ("2", "Metode deploy (Docker / Bee-Box VM)")]
cx = Inches(0.6); cw = Inches(3.9); gap = Inches(0.15)
for i, (big, sub) in enumerate(cards):
    x = cx + (cw + gap) * i
    rect(s, x, Inches(3.05), cw, Inches(2.2), LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, Inches(3.05), Inches(0.12), Inches(2.2), ACCENT)
    txt(s, x, Inches(3.35), cw, Inches(1.0),
        [[R(big, 52, NAVY, True, False, FONT)]], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(4.45), cw - Inches(0.4), Inches(0.7),
        [[R(sub, 14, TEXT, False, False, FONT)]], align=PP_ALIGN.CENTER)
txt(s, Inches(0.6), Inches(5.55), Inches(12), Inches(0.6),
    [[R("Catatan: bWAPP bukan aplikasi yang error, melainkan dirancang untuk "
        "dieksploitasi sebagai media latihan.", 13, GRAY, False, True, FONT)]])
footer(s)


# ================================================================ SLIDE 3: kenapa bWAPP (perbandingan)
s = slide()
header(s, "02  RELEVANSI", "Kenapa bWAPP untuk IT Audit?", "3")
bullets(s, Inches(0.6), Inches(1.45), Inches(5.7), Inches(4),
        [("Coverage terluas: ", "100+ kerentanan dalam satu platform."),
         ("Selaras OWASP: ", "tiap bug mudah dipetakan ke framework standar."),
         ("Setup cepat: ", "siap pakai via Docker dalam hitungan menit."),
         ("Bertingkat: ", "3 level kesulitan untuk simulasi bertahap.")],
        size=15, gap=12)
# comparison table
data = [["Platform", "Jumlah Bug"],
        ["bWAPP", "100+"],
        ["Juice Shop", "~75"],
        ["WebGoat", "~35"],
        ["DVWA", "~15"]]
gt = table(s, Inches(6.7), Inches(1.55), Inches(6.0), len(data), 2,
           [Inches(4.0), Inches(2.0)], row_height=Inches(0.62))
for i, row in enumerate(data):
    head = (i == 0)
    hl = (row[0] == "bWAPP")
    for j, val in enumerate(row):
        fill = NAVY if head else (RGBColor(0xDD, 0xEB, 0xF6) if hl else (LIGHT if i % 2 else WHITE))
        col = WHITE if head else (NAVY if hl else TEXT)
        set_cell(gt.cell(i, j), val, size=14, color=col, bold=(head or hl), fill=fill,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
txt(s, Inches(6.7), Inches(5.5), Inches(6), Inches(0.5),
    [[R("bWAPP mengungguli platform sejenis dalam cakupan kerentanan.",
        13, GRAY, False, True, FONT)]])
footer(s)


# ================================================================ SLIDE 4: setup & lab
s = slide()
header(s, "03  LINGKUNGAN", "Setup dan Lingkungan Lab", "4")
txt(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.4),
    [[R("Deployment via Docker (terisolasi ke localhost):", 15, NAVY2, True, False, FONT)]])
codebox(s, Inches(0.6), Inches(1.85), Inches(6.1), Inches(1.5),
        ["$ docker run -d -p 8080:80 \\",
         "    hackersploit/bwapp-docker",
         "",
         "# inisialisasi DB sekali:",
         "  http://localhost:8080/install.php"], size=12.5)
bullets(s, Inches(0.6), Inches(3.6), Inches(6.1), Inches(3),
        [("Akses: ", "http://localhost:8080"),
         ("Login: ", "bee / bug (kredensial lab default)"),
         ("Security level: ", "Low (security_level = 0)"),
         ("Stack: ", "PHP 5.5.9, Apache 2.4.7, MySQL, Ubuntu")],
        size=14, gap=10)
# right panel: scope summary
rect(s, Inches(7.05), Inches(1.85), Inches(5.7), Inches(4.55), PANEL,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.35), Inches(2.05), Inches(5.2), Inches(0.4),
    [[R("Ruang Lingkup Audit", 16, NAVY, True, False, FONT)]])
bullets(s, Inches(7.35), Inches(2.6), Inches(5.1), Inches(3.6),
        [("Jenis uji: ", "black-box dan grey-box"),
         ("In-scope: ", "web app layer (injection, XSS, info disclosure, misconfig, header)"),
         ("Out-of-scope: ", "audit OS/host, DoS, social engineering, sistem nyata"),
         ("Sifat: ", "lab terisolasi, aplikasi sengaja rentan")],
        size=13.5, gap=11)
footer(s)


# ================================================================ SLIDE 5: metodologi 5 tahap
s = slide()
header(s, "04  METODOLOGI", "Alur Audit: 5 Tahap", "5")
steps = [("1", "Setup", "Deploy bWAPP via Docker dan verifikasi akses."),
         ("2", "Scope", "Pilih kategori kerentanan dan tetapkan level kesulitan."),
         ("3", "Pengujian", "Scan otomatis, eksploitasi terarah, verifikasi manual."),
         ("4", "Analisis & Mapping", "Petakan ke OWASP Top 10 2021, tetapkan severity."),
         ("5", "Pelaporan", "Susun temuan, bukti, dan rekomendasi remediasi.")]
y = Inches(1.7); ch = Inches(0.92); gap = Inches(0.15)
for i, (n, title, desc) in enumerate(steps):
    yy = y + (ch + gap) * i
    rect(s, Inches(0.6), yy, Inches(12.1), ch, LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    circ = rect(s, Inches(0.85), yy + Inches(0.16), Inches(0.6), Inches(0.6),
                NAVY, shape=MSO_SHAPE.OVAL)
    ctf = circ.text_frame; ctf.margin_top = Pt(0); ctf.margin_bottom = Pt(0)
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = n; cr.font.size = Pt(22)
    cr.font.bold = True; cr.font.color.rgb = WHITE; cr.font.name = FONT
    txt(s, Inches(1.7), yy, Inches(3.0), ch,
        [[R(title, 18, NAVY, True, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.7), yy, Inches(7.8), ch,
        [[R(desc, 15, TEXT, False, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s)


# ================================================================ SLIDE 6: tools
s = slide()
header(s, "04  METODOLOGI", "Tools yang Digunakan", "6")
data = [["Tool", "Versi", "Fungsi"],
        ["Nikto", "2.1.5", "Web server scanner (info disclosure, misconfiguration)"],
        ["SQLMap", "1.10.2.6", "Deteksi dan eksploitasi SQL Injection otomatis"],
        ["OWASP ZAP", "2.17.0", "Dynamic Application Security Testing (DAST)"],
        ["Retire.js", "5.4.3", "Deteksi pustaka JavaScript usang (SCA)"],
        ["Nmap + vulners", "7.94SVN", "Deteksi versi layanan dan pemetaan CVE server"],
        ["Manual (curl, browser)", "-", "Verifikasi manual XSS dan kebocoran data"]]
gt = table(s, Inches(0.6), Inches(1.55), Inches(12.1), len(data), 3,
           [Inches(3.0), Inches(1.7), Inches(7.4)], row_height=Inches(0.68))
for i, row in enumerate(data):
    head = (i == 0)
    for j, val in enumerate(row):
        fill = NAVY if head else (LIGHT if i % 2 else WHITE)
        col = WHITE if head else TEXT
        set_cell(gt.cell(i, j), val, size=14, color=col, bold=head, fill=fill,
                 align=PP_ALIGN.CENTER if j == 1 else PP_ALIGN.LEFT)
txt(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
    [[R("Kombinasi scan otomatis + eksploitasi terarah + analisis komponen + "
        "uji manual untuk menutup blind spot tiap tool.", 13, GRAY, False, True, FONT)]])
footer(s)


# ================================================================ SLIDE 7: exec summary + severity
s = slide()
header(s, "05  HASIL", "Ringkasan Eksekutif", "7")
txt(s, Inches(0.6), Inches(1.45), Inches(7.0), Inches(3.6),
    [[R("Audit menemukan ", 16, TEXT, False, False, FONT),
      R("5 area temuan utama", 16, NAVY, True, False, FONT),
      R(" dengan severity Medium hingga Critical.", 16, TEXT, False, False, FONT)],
     [R("", 6, TEXT, False, False, FONT)],
     [R("Dua temuan Critical memungkinkan penyerang mengambil alih data "
        "sensitif tanpa autentikasi:", 15, TEXT, False, False, FONT)],
     [R("kebocoran kredensial", 15, CRIT, True, False, FONT),
      R(" melalui file konfigurasi yang terekspos, dan ", 15, TEXT, False, False, FONT),
      R("SQL Injection", 15, CRIT, True, False, FONT),
      R(" yang membongkar seluruh tabel pengguna termasuk akun admin.",
        15, TEXT, False, False, FONT)],
     [R("", 6, TEXT, False, False, FONT)],
     [R("Aplikasi juga rentan XSS, memakai komponen End-of-Life, dan memiliki "
        "sejumlah security misconfiguration.", 15, TEXT, False, False, FONT)]],
    line_spacing=1.12, space_after=6)
# severity distribution panel
rect(s, Inches(8.0), Inches(1.55), Inches(4.7), Inches(4.9), PANEL,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.0), Inches(1.75), Inches(4.7), Inches(0.4),
    [[R("Distribusi Severity", 16, NAVY, True, False, FONT)]],
    align=PP_ALIGN.CENTER)
dist = [("Critical", 2, CRIT), ("High", 2, HIGH), ("Medium", 1, MED)]
yy = Inches(2.45)
for label, cnt, color in dist:
    rect(s, Inches(8.35), yy, Inches(2.3), Inches(0.55), color,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(8.35), yy, Inches(2.3), Inches(0.55),
        [[R(label, 15, WHITE if color in (CRIT, HIGH) else TEXT, True, False, FONT)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(10.85), yy, Inches(1.6), Inches(0.55),
        [[R(str(cnt), 30, NAVY, True, False, FONT)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy = yy + Inches(0.85)
txt(s, Inches(8.0), Inches(5.35), Inches(4.7), Inches(0.9),
    [[R("Total 5 temuan utama", 16, NAVY, True, False, FONT)],
     [R("(plus 6 temuan Low pendukung)", 12, GRAY, False, False, FONT)]],
    align=PP_ALIGN.CENTER, space_after=2)
footer(s)


# ================================================================ SLIDE 8: ringkasan temuan (tabel)
s = slide()
header(s, "05  HASIL", "Ringkasan Temuan", "8")
data = [["ID", "Temuan", "Severity", "OWASP 2021", "Tool"],
        ["F-01", "Sensitive Data Exposure (kredensial bocor)", "Critical", "A01, A02, A05", "Nikto + manual"],
        ["F-02", "SQL Injection", "Critical", "A03, A02", "SQLMap"],
        ["F-03", "Cross-Site Scripting (Reflected + Stored)", "High", "A03", "Manual"],
        ["F-04", "Security Misconfiguration & Missing Headers", "Medium", "A05, A01", "OWASP ZAP"],
        ["F-05", "Vulnerable & Outdated Components", "High", "A06", "Retire.js + Nmap"]]
gt = table(s, Inches(0.6), Inches(1.55), Inches(12.1), len(data), 5,
           [Inches(0.9), Inches(5.0), Inches(1.7), Inches(2.5), Inches(2.0)],
           row_height=Inches(0.66))
for i, row in enumerate(data):
    head = (i == 0)
    for j, val in enumerate(row):
        if head:
            set_cell(gt.cell(i, j), val, size=13.5, color=WHITE, bold=True,
                     fill=NAVY, align=PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT)
        elif j == 2:
            set_cell(gt.cell(i, j), val, size=13, color=WHITE if val in ("Critical", "High") else TEXT,
                     bold=True, fill=sev_color(val), align=PP_ALIGN.CENTER)
        else:
            set_cell(gt.cell(i, j), val, size=13, color=TEXT,
                     fill=LIGHT if i % 2 else WHITE,
                     align=PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER)
footer(s)


# ================================================================ findings detail factory
def finding_slide(num, fid, title, sev, owasp, desc, evid_title, evid_lines,
                  impact, recs):
    s = slide()
    color = sev_color(sev)
    # header bar tinted with severity
    rect(s, 0, 0, EMU_W, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), EMU_W, Inches(0.06), color)
    txt(s, Inches(0.6), Inches(0.16), Inches(11), Inches(0.3),
        [[R("06  TEMUAN TERPERINCI", 12, RGBColor(0x9F, 0xC5, 0xE8), True, False, FONT)]])
    txt(s, Inches(0.6), Inches(0.40), Inches(9.5), Inches(0.7),
        [[R(f"{fid}   {title}", 24, WHITE, True, False, FONT)]],
        anchor=MSO_ANCHOR.MIDDLE)
    chip(s, Inches(11.2), Inches(0.38), sev.upper(), color, w=Inches(1.55), h=Inches(0.42))
    txt(s, Inches(12.4), Inches(0.16), Inches(0.7), Inches(0.3),
        [[R(num, 11, RGBColor(0x9F, 0xC5, 0xE8), False, False, FONT)]],
        align=PP_ALIGN.RIGHT)
    # OWASP line
    txt(s, Inches(0.6), Inches(1.32), Inches(12), Inches(0.35),
        [[R("OWASP 2021:  ", 13, NAVY2, True, False, FONT),
          R(owasp, 13, TEXT, False, False, FONT)]])
    # description
    txt(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.9),
        [[R(desc, 14.5, TEXT, False, False, FONT)]], line_spacing=1.08)
    # left: evidence
    txt(s, Inches(0.6), Inches(2.85), Inches(6.0), Inches(0.35),
        [[R(evid_title, 14, NAVY, True, False, FONT)]])
    codebox(s, Inches(0.6), Inches(3.25), Inches(6.0), Inches(2.4), evid_lines, size=11)
    # impact below code
    rect(s, Inches(0.6), Inches(5.8), Inches(6.0), Inches(1.05), RGBColor(0xFF, 0xF3, 0xCD),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.6), Inches(5.8), Inches(0.12), Inches(1.05), color)
    txt(s, Inches(0.85), Inches(5.9), Inches(5.6), Inches(0.9),
        [[R("Dampak:  ", 12.5, NAVY2, True, False, FONT),
          R(impact, 12.5, TEXT, False, False, FONT)]], line_spacing=1.05)
    # right: recommendations
    rect(s, Inches(6.9), Inches(2.85), Inches(5.8), Inches(4.0), PANEL,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(7.15), Inches(3.0), Inches(5.4), Inches(0.4),
        [[R("Rekomendasi Remediasi", 15, NAVY, True, False, FONT)]])
    bullets(s, Inches(7.15), Inches(3.5), Inches(5.35), Inches(3.2),
            recs, size=13, gap=9)
    footer(s)
    return s


# ---- F-01
finding_slide(
    "9", "F-01", "Sensitive Data Exposure", "Critical",
    "A01 Broken Access Control, A02 Cryptographic Failures, A05 Security Misconfiguration",
    "File konfigurasi dan backup berisi kredensial dapat diakses publik tanpa "
    "autentikasi. Directory indexing yang aktif menampilkan isi folder sensitif "
    "seperti /passwords/.",
    "Bukti: config.inc terekspos",
    ["$ curl http://localhost:8080/config.inc",
     "",
     '$server   = "localhost";',
     '$username = "bwapp";',
     '$password = "bwApped";',
     '$database = "bWAPP";',
     "",
     "/passwords/  -> heroes.xml, web.config.bak,",
     "               wp-config.bak (password plaintext)"],
    "kredensial DB diperoleh tanpa eksploitasi, berpotensi full database compromise.",
    ["Matikan directory indexing (Options -Indexes).",
     "Hapus file backup (.bak, .inc, .old) dari web root.",
     "Pindahkan file konfigurasi ke luar document root.",
     "Jangan simpan password plaintext, gunakan hashing kuat.",
     "Hapus installer (install.php) setelah instalasi."])

# ---- F-02
finding_slide(
    "10", "F-02", "SQL Injection", "Critical",
    "A03 Injection, A02 Cryptographic Failures",
    "Parameter title pada modul SQL Injection (GET/Search) tidak disanitasi, "
    "sehingga input pengguna dieksekusi langsung sebagai bagian dari query SQL.",
    "Bukti: SQLMap dump bWAPP.users",
    ["Parameter 'title' rentan 4 teknik:",
     "boolean-blind, error, time-blind, UNION",
     "",
     "Database: bWAPP  Table: users",
     "id | login  | admin | password (SHA1)",
     "1  | A.I.M. |  1    | 6885858486...affd0",
     "2  | bee    |  1    | 6885858486...affd0",
     "",
     'echo -n "bug" | sha1sum -> match',
     "Password admin dipulihkan: \"bug\""],
    "akses baca penuh ke database; hash SHA1 unsalted mudah dipecahkan.",
    ["Gunakan prepared statements / parameterized query (PDO, MySQLi).",
     "Terapkan input validation berbasis whitelist.",
     "Ganti hashing password ke bcrypt atau argon2 dengan salt unik.",
     "Terapkan least privilege pada akun database aplikasi."])

# ---- F-03
finding_slide(
    "11", "F-03", "Cross-Site Scripting (XSS)", "High",
    "A03 Injection (Cross-Site Scripting)",
    "Aplikasi tidak melakukan output encoding maupun sanitasi input. Ditemukan "
    "dua varian: Reflected (GET) dan Stored (Blog) yang tersimpan permanen.",
    "Bukti: payload tereksekusi",
    ["Reflected (xss_get.php, param firstname):",
     "Welcome <script>alert('XSS-bWAPP')",
     "        </script> test",
     "",
     "Stored (xss_stored_1.php):",
     "<td><script>alert('Stored-XSS-by-",
     "     Audit')</script></td>",
     "",
     "(payload stored telah dibersihkan",
     " setelah verifikasi)"],
    "pencurian session cookie, session hijacking, defacement, phishing. "
    "Diperparah cookie tanpa HttpOnly.",
    ["Terapkan output encoding sesuai konteks (htmlspecialchars).",
     "Lakukan input validation di sisi server.",
     "Pasang Content Security Policy (CSP).",
     "Set flag HttpOnly dan Secure pada cookie session."])

# ---- F-04
finding_slide(
    "12", "F-04", "Security Misconfiguration", "Medium",
    "A05 Security Misconfiguration, A01 Broken Access Control",
    "OWASP ZAP menemukan sejumlah kesalahan konfigurasi dan header keamanan yang "
    "hilang: total 5 temuan Medium dan 6 temuan Low pendukung.",
    "Bukti: temuan Medium dari ZAP",
    ["- Absence of Anti-CSRF Tokens   (CWE-352)",
     "- Application Error Disclosure  (CWE-550)",
     "- CSP Header Not Set            (CWE-693)",
     "- Directory Browsing            (CWE-548)",
     "- Missing Anti-clickjacking     (CWE-1021)",
     "",
     "Low: cookie tanpa HttpOnly/SameSite,",
     "banner Apache/2.4.7 & PHP/5.5.9 bocor,",
     "X-Content-Type-Options tak diset"],
    "memperbesar peluang CSRF, clickjacking, dan info disclosure bagi penyerang.",
    ["Tambahkan token Anti-CSRF unik pada setiap form pengubah state.",
     "Set header CSP, X-Frame-Options, X-Content-Type-Options: nosniff.",
     "Set flag HttpOnly, Secure, SameSite pada cookie.",
     "Gunakan custom error page; sembunyikan banner versi server.",
     "Update komponen EOL (PHP 5.5.9, Apache 2.4.7)."])

# ---- F-05
finding_slide(
    "13", "F-05", "Vulnerable & Outdated Components", "High",
    "A06 Vulnerable and Outdated Components",
    "Aplikasi berjalan di atas komponen End-of-Life pada sisi klien (jQuery) dan "
    "server (Apache, PHP) yang tidak lagi menerima patch keamanan.",
    "Bukti: Retire.js + Nmap vulners",
    ["jQuery 1.4.4 (2010, EOL) -> 7 kerentanan:",
     "CVE-2020-11023, CVE-2019-11358,",
     "CVE-2020-7656, CVE-2015-9251, dll",
     "",
     "Apache httpd 2.4.7 (Ubuntu):",
     "vulners cocokkan 106 CVE unik",
     "(Critical 24 | High 42 | Medium 40)",
     "",
     "PHP 5.5.9 -> EOL sejak Juli 2016"],
    "permukaan serangan luas (XSS, prototype pollution, request smuggling) dengan "
    "banyak exploit publik.",
    ["Perbarui jQuery ke versi didukung (3.5.0+).",
     "Update Apache; migrasikan PHP 5.5.9 ke PHP 8.x.",
     "Terapkan patch management dan SCA berkala.",
     "Catatan: skor per-CVE bersifat potensial (berbasis banner); "
     "fakta EOL yang terkonfirmasi bernilai High."])


# ================================================================ SLIDE 14: catatan metodologi
s = slide()
header(s, "07  CATATAN METODOLOGI", "Keterbatasan Automated Scanner", "14")
# highlight box
rect(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(1.25), RGBColor(0xFF, 0xF3, 0xCD),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(1.45), Inches(0.14), Inches(1.25), HIGH)
txt(s, Inches(0.95), Inches(1.6), Inches(11.5), Inches(1.0),
    [[R("OWASP ZAP melaporkan 0 temuan High", 18, CRIT, True, False, FONT),
      R(", padahal pengujian manual dan SQLMap membuktikan adanya SQL Injection "
        "(Critical) dan Stored XSS (High) pada aplikasi yang sama.",
        16, TEXT, False, False, FONT)]], line_spacing=1.1,
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.6), Inches(2.95), Inches(12.1), Inches(1.1),
    [[R("Penyebab: ZAP hanya meng-crawl halaman permukaan dan tidak masuk ke "
        "modul kerentanan di balik autentikasi atau yang tidak ter-link otomatis. "
        "Automated DAST scanner memiliki blind spot, sehingga audit yang baik wajib "
        "mengombinasikan scan otomatis, eksploitasi terarah, dan uji manual.",
        15, TEXT, False, False, FONT)]], line_spacing=1.12)
txt(s, Inches(0.6), Inches(4.25), Inches(12), Inches(0.4),
    [[R("Cross-Validation (dikonfirmasi 2 tool berbeda: Nikto + ZAP)", 16, NAVY, True, False, FONT)]])
bullets(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(2),
        ["Directory Browsing aktif",
         "Kebocoran banner versi server (Apache/2.4.7)",
         "Missing X-Frame-Options (anti-clickjacking)",
         "Cookie session tanpa flag HttpOnly"], size=15, gap=8)
footer(s)


# ================================================================ SLIDE 15: best practice / etika
s = slide()
header(s, "08  ETIKA LAB", "Best Practice dan Keamanan Lab", "15")
# DO panel
rect(s, Inches(0.6), Inches(1.55), Inches(5.9), Inches(4.9),
     RGBColor(0xE8, 0xF3, 0xE8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(0.5),
    [[R("YANG HARUS DILAKUKAN", 17, RGBColor(0x1E, 0x7A, 0x1E), True, False, FONT)]],
    align=PP_ALIGN.CENTER)
bullets(s, Inches(0.95), Inches(2.4), Inches(5.3), Inches(4),
        ["Bind Docker hanya ke localhost.",
         "Gunakan Host-only Adapter pada VM.",
         "Ambil snapshot sebelum eksploitasi.",
         "Dokumentasikan setiap langkah (payload, screenshot, respons).",
         "Bersihkan payload stored setelah verifikasi."],
        size=15, gap=12, color=TEXT)
# DON'T panel
rect(s, Inches(6.8), Inches(1.55), Inches(5.9), Inches(4.9),
     RGBColor(0xFB, 0xE9, 0xE9), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(6.8), Inches(1.7), Inches(5.9), Inches(0.5),
    [[R("YANG TIDAK BOLEH", 17, CRIT, True, False, FONT)]],
    align=PP_ALIGN.CENTER)
bullets(s, Inches(7.15), Inches(2.4), Inches(5.3), Inches(4),
        ["Expose port ke jaringan publik.",
         "Menjalankan di komputer produksi.",
         "Menggunakan payload terhadap target nyata tanpa izin tertulis (ilegal).",
         "Memperlakukan data dump sebagai kredensial nyata (ini user seeded bWAPP)."],
        size=15, gap=12, color=TEXT)
footer(s)


# ================================================================ SLIDE 16: kesimpulan
s = slide()
header(s, "09  KESIMPULAN", "Kesimpulan", "16")
bullets(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(4.5),
        [("Kerentanan serius dan beragam. ",
          "bWAPP mencakup mayoritas kategori OWASP Top 10 2021."),
         ("Dua Critical jadi prioritas. ",
          "Sensitive Data Exposure dan SQL Injection memungkinkan kompromi data langsung."),
         ("Komponen usang memperluas serangan. ",
          "A06: jQuery 1.4.4, Apache 2.4.7, dan PHP 5.5.9 sudah End-of-Life."),
         ("Kombinasi metode terbukti penting. ",
          "Tidak ada satu tool pun yang menemukan seluruh kerentanan."),
         ("Media latihan audit yang efektif. ",
          "Dari identifikasi, eksploitasi, hingga rekomendasi selaras standar industri.")],
        size=17, gap=18)
footer(s)


# ================================================================ SLIDE 17: penutup
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, NAVY)
rect(s, 0, Inches(4.6), EMU_W, Inches(0.07), ACCENT)
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2),
    [[R("Terima Kasih", 48, WHITE, True, False, FONT)]])
txt(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.6),
    [[R("Ada pertanyaan?", 22, RGBColor(0xCF, 0xE0, 0xF0), False, False, FONT)]])
txt(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.8),
    [[R("Auditor:  ", 16, GRAY, True, False, FONT),
      R("Ravi Arnan Irianto  -  Mata Kuliah IT Audit", 16, WHITE, False, False, FONT)],
     [R("Demo cepat:  ", 16, GRAY, True, False, FONT),
      R("http://localhost:8080   (login bee / bug)", 16, WHITE, False, False, FONT)],
     [R("Referensi:  ", 16, GRAY, True, False, FONT),
      R("owasp.org/Top10  -  itsecgames.com  -  cheatsheetseries.owasp.org",
        16, WHITE, False, False, FONT)]],
    space_after=8)

import os
os.makedirs("doc", exist_ok=True)
out = "doc/presentasi_audit_bwapp.pptx"
prs.save(out)
print(f"saved {out} with {len(prs.slides._sldIdLst)} slides")
